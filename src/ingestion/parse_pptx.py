import os
import json
from pptx import Presentation
from PIL import Image
import pytesseract

PPTX_PATH = "data/raw/2023 Reporting Portal Trunk- Manager Guide .pptx"
IMAGE_DIR = "data/images/pptx"
OUTPUT_JSON = "data/processed/pptx_chunks.json"

os.makedirs(IMAGE_DIR, exist_ok=True)

prs = Presentation(PPTX_PATH)
all_slides_content = []

for slide_idx, slide in enumerate(prs.slides):
    slide_number = slide_idx + 1  # Always use the correct slide number!
    # Sort shapes by (top, left)
    sorted_shapes = sorted(
        slide.shapes, 
        key=lambda s: (
            getattr(s, "top", 0), 
            getattr(s, "left", 0)
        )
    )
    slide_content = []
    for shape in sorted_shapes:
        # If text
        if hasattr(shape, "text") and shape.text.strip():
            slide_content.append({"type": "text", "content": shape.text.strip()})
        # If image
        if hasattr(shape, "image"):
            try:
                image = shape.image
                image_ext = image.ext
                img_filename = f"{IMAGE_DIR}/slide{slide_number}_img{shape.shape_id}.{image_ext}"
                with open(img_filename, "wb") as img_out:
                    img_out.write(image.blob)
                # OCR image
                try:
                    img = Image.open(img_filename)
                    ocr_text = pytesseract.image_to_string(img).strip()
                    slide_content.append({
                        "type": "image",
                        "file": img_filename,
                        "ocr_text": ocr_text
                    })
                except Exception as e:
                    slide_content.append({
                        "type": "image",
                        "file": img_filename,
                        "ocr_text": f"OCR error: {e}"
                    })
            except Exception:
                continue
    # MAIN FIX: Each chunk now gets proper source/citation detail
    all_slides_content.append({
        "chunk_id": slide_number,
        "source": "pptx",
        "source_detail": f"slide {slide_number}",
        "text": "\n".join(block["content"] for block in slide_content if block["type"] == "text"),
        "blocks": slide_content,
        "metadata": {"slide_number": slide_number}
    })

# Print a summary for inspection
for slide in all_slides_content:
    print(f"\n{'='*20}\nSlide {slide['metadata']['slide_number']}")
    for block in slide["blocks"]:
        if block["type"] == "text":
            print(f"[Text] {block['content']}")
        else:
            print(f"[Image OCR] {block['ocr_text'][:120]}{'...' if len(block['ocr_text'])>120 else ''}")

# Save to JSON
with open(OUTPUT_JSON, "w") as f:
    json.dump(all_slides_content, f, indent=2)

print(f"\nPPTX slides parsed and chunked. Output saved to {OUTPUT_JSON}")
