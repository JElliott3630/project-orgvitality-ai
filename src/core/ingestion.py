from __future__ import annotations
"""DocumentIngestor – single-responsibility class
--------------------------------------------------
• Supports .pdf, .docx, .xlsx, .pptx (easily extensible)
• Converts each file to markdown-ish plain‑text "pages"
• Splits pages into token windows (750 tokens, 150 overlap)
• Returns a DocumentBatch dataclass ready for embedding / storage

Public API
~~~~~~~~~~
    ingest(file_bytes: bytes, file_name: str)  -> DocumentBatch
    ingest_path(path: str | Path)             -> DocumentBatch
"""

import unicodedata
import uuid
from io import BytesIO
from pathlib import Path
from typing import List

import pdfplumber
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import tiktoken
from pptx import Presentation
from docx2pdf import convert as docx2pdf_convert
import tempfile

from core.schema import Chunk, DocumentBatch

__all__ = ["DocumentIngestor"]

ALLOWED_EXT = {".pdf", ".docx", ".xlsx", ".pptx"}


class DocumentIngestor:
    """Convert heterogeneous files to token‑chunks."""

    def __init__(self, chunk_size: int = 500, overlap: int = 75):
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.encoder = tiktoken.get_encoding("cl100k_base")

    # ------------------------------------------------------------------
    # Public helpers
    # ------------------------------------------------------------------

    def ingest(self, file_bytes: bytes, file_name: str) -> DocumentBatch:  # noqa: D401
        """Ingest an in‑memory file.*"""
        ext = self._detect_extension(file_name)
        pages = self._parse_to_pages(file_bytes, ext)
        return self._pages_to_batch(pages, file_name)

    def ingest_path(self, path: str | Path) -> DocumentBatch:
        path = Path(path)
        with path.open("rb") as f:
            return self.ingest(f.read(), path.name)

    # ------------------------------------------------------------------
    # Private — detection / parsing
    # ------------------------------------------------------------------

    @staticmethod
    def _detect_extension(file_name: str) -> str:
        ext = Path(file_name).suffix.lower()
        if ext not in ALLOWED_EXT:
            raise ValueError(f"Unsupported file extension: {ext}")
        return ext

    def _parse_to_pages(self, file_bytes: bytes, ext: str) -> List[str]:
        match ext:
            case ".pdf":
                return self._parse_pdf(file_bytes)
            case ".docx":
                return self._parse_docx(file_bytes)
            case ".xlsx":
                return self._parse_excel(file_bytes)
            case ".pptx":
                return self._parse_pptx(file_bytes)
            case _:
                raise AssertionError("unreachable")

    # ---- format‑specific parsers -------------------------------------

    def _parse_pdf(self, file_bytes: bytes) -> List[str]:
        pages = []
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                text = unicodedata.normalize("NFKC", text).strip()
                if text:
                    pages.append(text)
        return pages or [""]

    def _parse_docx(self, file_bytes: bytes) -> List[str]:
        doc = DocxDocument(BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return ["\n".join(paragraphs)]

    def _parse_excel(self, file_bytes: bytes) -> List[str]:
        xl = pd.ExcelFile(BytesIO(file_bytes))
        pages = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            pages.append(f"### Sheet: {sheet}\n" + df.to_markdown(index=False))
        return pages or [""]

    def _parse_pptx(self, file_bytes: bytes) -> List[str]:
        prs = Presentation(BytesIO(file_bytes))
        pages = []
        for slide in prs.slides:
            texts = [s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
            pages.append("\n".join(texts))
        return pages or [""]

    # ---- chunking -----------------------------------------------------

    def _pages_to_batch(self, pages: List[str], source_name: str) -> DocumentBatch:
        doc_id = str(uuid.uuid4())
        chunks: List[Chunk] = []
        for p_idx, page_text in enumerate(pages, 1):
            location = f"page {p_idx}" if len(pages) > 1 else "full"
            tokens = self.encoder.encode(page_text)
            for i in range(0, len(tokens), self.chunk_size - self.overlap):
                window = tokens[i : i + self.chunk_size]
                chunk_text = self.encoder.decode(window)
                chunk_id = f"{doc_id}_{p_idx}_{i}"
                chunks.append(
                    Chunk(
                        id=chunk_id,
                        text=chunk_text,
                        source=source_name,
                        location=location,
                        start_token=i,
                    )
                )
        return DocumentBatch(doc_id=doc_id, chunks=chunks)
